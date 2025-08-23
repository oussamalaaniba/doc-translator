import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def iter_all_shapes(shapes):
    """Itère sur toutes les formes, y compris à l'intérieur des GroupShapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all_shapes(shape.shapes)
        else:
            yield shape

def _paragraph_text(paragraph):
    """Texte du paragraphe avec les espaces corrects (via .text)."""
    return paragraph.text or ""

def _replace_paragraph_text_keep_font(paragraph, new_text):
    """
    Remplace le contenu d'un paragraphe SANS supprimer les runs :
    - met tout le texte dans le 1er run (il garde police/taille),
    - vide les runs suivants proprement (r.text = ""),
    - crée 1 run si le paragraphe était vide.
    """
    if paragraph.runs:
        first = paragraph.runs[0]
        first.text = new_text           # conserve police/size
        for r in paragraph.runs[1:]:
            r.text = ""                 # nettoie sans toucher au XML
    else:
        run = paragraph.add_run()
        run.text = new_text

def _translate_text_frame_by_paragraph(text_frame, translate_callable, src, tgt):
    """Traduit chaque PARAGRAPHE (pas run) pour éviter pertes d'espaces et mot-à-mot."""
    # collecter d'abord -> batch unique -> remplacer
    para_refs, texts = [], []
    for p in list(text_frame.paragraphs):
        original = _paragraph_text(p)
        if original and original.strip():
            para_refs.append(p)
            texts.append(original)

    if not para_refs:
        return

    translated_list = translate_callable(texts, src, tgt) if translate_callable else texts
    for p, new_text in zip(para_refs, translated_list):
        _replace_paragraph_text_keep_font(p, new_text)

def _translate_table_by_paragraph(table, translate_callable, src, tgt):
    for row in table.rows:
        for cell in row.cells:
            if hasattr(cell, "text_frame") and cell.text_frame:
                _translate_text_frame_by_paragraph(cell.text_frame, translate_callable, src, tgt)

def _translate_chart_titles_if_any(shape, translate_callable, src, tgt):
    """Traduit titres de graphiques et titres d’axes si exposés par python-pptx."""
    try:
        if hasattr(shape, "has_chart") and shape.has_chart:
            chart = shape.chart
            # Titre du chart
            if chart.has_title and hasattr(chart.chart_title, "text_frame"):
                _translate_text_frame_by_paragraph(chart.chart_title.text_frame, translate_callable, src, tgt)
            # Titres d’axes (selon versions)
            for axis_name in ("category_axis", "value_axis", "series_axis"):
                axis = getattr(chart, axis_name, None)
                if axis and getattr(axis, "has_title", False) and hasattr(axis.axis_title, "text_frame"):
                    _translate_text_frame_by_paragraph(axis.axis_title.text_frame, translate_callable, src, tgt)
    except Exception:
        # selon la version de python-pptx, certaines propriétés peuvent manquer
        pass

def translate_pptx_preserve_styles(src_bytes, src="fr", tgt="en", translate_callable=None):
    """
    Traduction PAR PARAGRAPHE (évite perte d'espaces et mot-à-mot) avec styles conservés :
    - Zones de texte / placeholders
    - Objets groupés (récursif)
    - Tableaux
    - Titres de graphiques (+ axes si exposés)
    - SmartArt/images avec texte: non supportés
    """
    prs = Presentation(io.BytesIO(src_bytes))

    for slide in prs.slides:
        for shape in iter_all_shapes(slide.shapes):
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                _translate_text_frame_by_paragraph(shape.text_frame, translate_callable, src, tgt)

            if getattr(shape, "has_table", False) or shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                _translate_table_by_paragraph(shape.table, translate_callable, src, tgt)

            _translate_chart_titles_if_any(shape, translate_callable, src, tgt)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()
